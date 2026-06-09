from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_PATH = BASE_DIR / "專題簡報_電商網站專案_初稿.pptx"

BG = RGBColor(246, 243, 238)
NAVY = RGBColor(47, 62, 107)
SLATE = RGBColor(102, 117, 163)
LIGHT = RGBColor(219, 225, 242)
TEXT = RGBColor(46, 46, 46)
MUTED = RGBColor(110, 110, 110)
ACCENT = RGBColor(175, 37, 65)
LINE = RGBColor(176, 188, 220)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_rule(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.78), Inches(12.1), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide, page_no):
    box = slide.shapes.add_textbox(Inches(5.55), Inches(7.0), Inches(1.2), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(page_no)
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_title(slide, title, subtitle=None):
    add_bg(slide)
    add_top_rule(slide)
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.2), Inches(11.0), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft JhengHei"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(2.2), Inches(0.58), Inches(7.6), Inches(0.24))
        p = sub_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(10.5)
        run.font.color.rgb = MUTED


def add_section_slide(prs, title, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide)
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.0), Inches(1.4))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft JhengHei"
    run.font.bold = True
    run.font.size = Pt(30)
    run.font.color.rgb = NAVY
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(3.6), Inches(2.0), Inches(0.55))
    tag.fill.solid()
    tag.fill.fore_color.rgb = SLATE
    tag.line.color.rgb = RGBColor(255, 255, 255)
    p = tag.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "專題簡報"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    add_footer(slide, page_no)
    return slide


def add_bullet_slide(prs, title, bullets, page_no, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    left = slide.shapes.add_textbox(Inches(0.95), Inches(1.15), Inches(10.4), Inches(5.5))
    tf = left.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        run = p.add_run()
        run.text = bullet
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(20)
        run.font.color.rgb = TEXT
    add_footer(slide, page_no)
    return slide


def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items, page_no, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    for x, heading, items in [
        (0.7, left_title, left_items),
        (6.35, right_title, right_items),
    ]:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.15), Inches(5.0), Inches(5.35))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(252, 252, 252)
        card.line.color.rgb = LINE
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(1.35), Inches(1.9), Inches(0.42))
        head.fill.solid()
        head.fill.fore_color.rgb = SLATE
        head.line.color.rgb = SLATE
        p = head.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = heading
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        box = slide.shapes.add_textbox(Inches(x + 0.28), Inches(1.92), Inches(4.4), Inches(4.2))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = item
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(18)
            run.font.color.rgb = TEXT
    add_footer(slide, page_no)
    return slide


def add_system_architecture(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "系統架構", "以前後端分離與 MySQL 單一資料來源為核心")

    layers = [
        ("使用者端", 0.9, 1.35, 2.1, 0.62, SLATE),
        ("Next.js / React 前端", 3.15, 1.35, 2.6, 0.62, NAVY),
        ("Django / DRF API", 6.0, 1.35, 2.3, 0.62, NAVY),
        ("MySQL / Aiven", 8.55, 1.35, 2.1, 0.62, SLATE),
    ]
    for text, x, y, w, h, color in layers:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Microsoft JhengHei"
        run.font.bold = True
        run.font.size = Pt(17)
        run.font.color.rgb = RGBColor(255, 255, 255)

    for x1, x2 in [(3.0, 3.15), (5.75, 6.0), (8.3, 8.55)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(1.66), Inches(x2), Inches(1.66))
        line.line.color.rgb = LINE
        line.line.width = Pt(2)

    blocks = [
        ("頁面模組", ["首頁", "商品列表", "購物車", "會員中心", "賣家中心", "後台管理"], 0.85),
        ("後端服務", ["認證", "商品管理", "訂單流程", "金流物流", "社群內容", "Banner 與媒體"], 4.25),
        ("外部服務", ["Render 部署", "Google Cloud Storage", "NewebPay 金流", "Aiven MySQL"], 7.75),
    ]
    for heading, items, x in blocks:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.0), Inches(3.9))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(252, 252, 252)
        card.line.color.rgb = LINE
        hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(2.38), Inches(1.4), Inches(0.4))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = LIGHT
        hdr.line.color.rgb = LIGHT
        p = hdr.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = heading
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = NAVY
        box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.95), Inches(2.45), Inches(2.9))
        tf = box.text_frame
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(7)
            run = p.add_run()
            run.text = item
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(16)
            run.font.color.rgb = TEXT
    add_footer(slide, page_no)
    return slide


def add_db_slide(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "資料庫規劃", "以 Django ORM 管理 MySQL 結構，資料來源不再混用 JSON")
    groups = [
        ("會員資料", ["users", "user_addresses", "user_invoice_profiles", "user_shipping_rules", "seller_requests"], 0.65, 1.3),
        ("商品資料", ["products", "brands", "categories", "product_images", "product_variants", "tags"], 4.2, 1.3),
        ("交易資料", ["carts", "cart_items", "orders", "order_items", "payment_transactions", "shipment_events"], 7.8, 1.3),
        ("內容資料", ["product_reviews", "product_questions", "community_posts", "community_replies", "banners", "media_assets"], 2.45, 4.0),
        ("營運資料", ["admin_audit_logs", "payment_callback_logs", "recent_views", "compare_items", "favorites / recommendations"], 6.05, 4.0),
    ]
    for title, items, x, y in groups:
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(1.85))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(252, 252, 252)
        card.line.color.rgb = LINE
        hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y + 0.15), Inches(1.35), Inches(0.36))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = SLATE
        hdr.line.color.rgb = SLATE
        p = hdr.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.6), Inches(2.5), Inches(1.0))
        tf = box.text_frame
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(3)
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = "Consolas"
            run.font.size = Pt(12)
            run.font.color.rgb = TEXT
    note = slide.shapes.add_textbox(Inches(0.85), Inches(6.25), Inches(10.4), Inches(0.35))
    p = note.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "目前本地同步資料：users=8、products=11、orders=12、banners=2"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NAVY
    add_footer(slide, page_no)
    return slide


def add_flow_slide(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "使用者購物流程", "依照老師文件要求，以流程圖方式呈現主要操作路徑")
    nodes = [
        ("進入首頁", 0.8, 2.3, 1.4, 0.62, MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR),
        ("瀏覽商品", 2.45, 2.3, 1.45, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("查看詳情", 4.15, 2.3, 1.45, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("加入購物車", 5.85, 2.3, 1.55, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("會員登入?", 7.7, 2.25, 1.35, 0.72, MSO_AUTO_SHAPE_TYPE.DIAMOND),
        ("前往結帳", 9.35, 2.3, 1.35, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("建立訂單", 4.45, 4.45, 1.55, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("金流 / 物流", 6.45, 4.45, 1.55, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("訂單查詢", 8.45, 4.45, 1.45, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
        ("完成", 10.25, 4.45, 1.0, 0.62, MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR),
        ("導向登入", 7.55, 5.65, 1.65, 0.62, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE),
    ]
    for text, x, y, w, h, shape_type in nodes:
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SLATE if shape_type != MSO_AUTO_SHAPE_TYPE.DIAMOND else LIGHT
        shape.line.color.rgb = RGBColor(255, 255, 255) if shape_type != MSO_AUTO_SHAPE_TYPE.DIAMOND else SLATE
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255) if shape_type != MSO_AUTO_SHAPE_TYPE.DIAMOND else NAVY

    connectors = [
        (2.2, 2.61, 2.45, 2.61, ""),
        (3.9, 2.61, 4.15, 2.61, ""),
        (5.6, 2.61, 5.85, 2.61, ""),
        (7.4, 2.61, 7.7, 2.61, ""),
        (9.05, 2.61, 9.35, 2.61, "是"),
        (10.02, 2.92, 10.02, 4.45, ""),
        (9.8, 4.76, 10.25, 4.76, ""),
        (8.0, 4.76, 8.45, 4.76, ""),
        (6.0, 4.76, 6.45, 4.76, ""),
        (7.7, 2.97, 8.35, 5.65, "否"),
        (8.35, 5.65, 8.35, 5.65, ""),
        (8.35, 5.65, 8.35, 5.65, ""),
    ]
    for x1, y1, x2, y2, label in connectors:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = LINE
        line.line.width = Pt(2)
        if label:
            box = slide.shapes.add_textbox(Inches((x1 + x2) / 2), Inches((y1 + y2) / 2 - 0.18), Inches(0.35), Inches(0.2))
            p = box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = ACCENT

    back_line_1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.2), Inches(5.96), Inches(9.2), Inches(6.2))
    back_line_1.line.color.rgb = LINE
    back_line_1.line.width = Pt(2)
    back_line_2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.2), Inches(6.2), Inches(5.25), Inches(6.2))
    back_line_2.line.color.rgb = LINE
    back_line_2.line.width = Pt(2)
    back_line_3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.25), Inches(6.2), Inches(5.25), Inches(5.07))
    back_line_3.line.color.rgb = LINE
    back_line_3.line.width = Pt(2)
    caption = slide.shapes.add_textbox(Inches(0.95), Inches(1.2), Inches(10.2), Inches(0.45))
    p = caption.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "主流程涵蓋商品瀏覽、會員驗證、訂單建立與金流物流串接"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(15)
    run.font.color.rgb = TEXT
    add_footer(slide, page_no)
    return slide


def add_sitemap_slide(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "網站架構", "依前端頁面與 Django 路由整理的主要功能層級")
    groups = [
        ("會員驗證", ["登入", "註冊", "忘記密碼", "重設密碼"], 0.55, 1.4),
        ("商品瀏覽", ["商品列表", "商品詳情", "品牌頁", "分類頁", "商品比較"], 3.55, 1.4),
        ("購物流程", ["購物車", "結帳", "訂單列表", "訂單詳情"], 6.55, 1.4),
        ("會員中心", ["儀表板", "個人資料", "地址", "發票", "優惠資訊"], 9.55, 1.4),
        ("賣家中心", ["我的商品", "新增商品", "物流規則", "銷售訂單", "報表"], 2.05, 4.0),
        ("社群論壇", ["文章列表", "文章詳情", "留言互動"], 5.25, 4.0),
        ("管理後台", ["後台儀表板", "商品管理", "訂單審核", "使用者管理", "Banner / 內容審核"], 8.2, 4.0),
    ]
    for title, items, x, y in groups:
        width = 2.55 if y < 3 else 2.75
        height = 1.75 if len(items) <= 4 else 1.95
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(252, 252, 252)
        card.line.color.rgb = LINE
        hdr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y + 0.14), Inches(1.35), Inches(0.35))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = SLATE
        hdr.line.color.rgb = SLATE
        p = hdr.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.58), Inches(width - 0.28), Inches(height - 0.4))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT
    add_footer(slide, page_no)
    return slide


def add_results_slide(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "實驗結果", "以目前專案完成度整理可展示成果")

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.25), Inches(5.2), Inches(2.75))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(252, 252, 252)
    left.line.color.rgb = LINE
    box = slide.shapes.add_textbox(Inches(0.92), Inches(1.48), Inches(4.6), Inches(2.15))
    tf = box.text_frame
    results = [
        "前端已完成商品、會員、購物車、結帳、訂單、社群與後台頁面分流。",
        "後端以 Django ORM + MySQL 管理正式資料，避免 JSON 與資料庫混用。",
        "已串接 Google Cloud Storage 儲存媒體檔案。",
        "已整合 Render 部署設定與 NewebPay 金流 / 物流流程。",
    ]
    for idx, item in enumerate(results):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = item
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(18)
        run.font.color.rgb = TEXT

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(1.25), Inches(5.0), Inches(2.75))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(252, 252, 252)
    right.line.color.rgb = LINE
    stats = [
        ("同步狀態", "本地資料庫已與 Aiven 匯出資料同步"),
        ("目前資料", "users=8 / products=11 / orders=12 / banners=2"),
        ("部署型態", "Django 與 Next.js 分開部署"),
        ("程式分層", "services / api / frontend 三層結構"),
    ]
    y = 1.5
    for label, value in stats:
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(y), Inches(1.25), Inches(0.35))
        tag.fill.solid()
        tag.fill.fore_color.rgb = LIGHT
        tag.line.color.rgb = LIGHT
        p = tag.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(11.5)
        run.font.bold = True
        run.font.color.rgb = NAVY
        txt = slide.shapes.add_textbox(Inches(7.8), Inches(y - 0.02), Inches(2.95), Inches(0.38))
        p = txt.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(13.5)
        run.font.color.rgb = TEXT
        y += 0.56

    image_paths = [
        BASE_DIR / "static" / "images" / "abc-短袖上衣1.png",
        BASE_DIR / "static" / "images" / "abc-短褲1.png",
        BASE_DIR / "static" / "images" / "abc2-長袖上衣(白).png",
        BASE_DIR / "static" / "images" / "abc3-NEW FORCE-banner.jpg",
    ]
    placements = [
        (0.9, 4.4, 2.2, 1.6),
        (3.35, 4.4, 2.2, 1.6),
        (5.95, 4.4, 2.2, 1.6),
        (8.55, 4.4, 2.2, 1.6),
    ]
    for path, (x, y, w, h) in zip(image_paths, placements):
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    add_footer(slide, page_no)
    return slide


def add_future_slide(prs, page_no):
    bullets_left = [
        "補齊自動化測試，涵蓋認證、購物流程、訂單與後台 API。",
        "將賣家營運報表、審核流程與通知機制做成更完整的工作流。",
        "優化前端 SEO、圖片載入與快取策略，提升使用者體驗。",
    ]
    bullets_right = [
        "加強權限控管、審計記錄與敏感設定管理，貼近正式環境需求。",
        "補上更多資料分析，例如熱門商品、回購率與內容互動指標。",
        "未來可擴充優惠券、推薦系統與多種付款配送方案。",
    ]
    return add_two_col_slide(prs, "未來方向", "系統強化", bullets_left, "營運延伸", bullets_right, page_no)


def add_reference_slide(prs, page_no):
    bullets = [
        "課程文件：《專題說明(20240822)》",
        "範例簡報：《購物網站架設.pptx》",
        "Django 官方文件：https://docs.djangoproject.com/",
        "Next.js 官方文件：https://nextjs.org/docs",
        "Render 官方文件：https://render.com/docs",
        "Aiven for MySQL 文件：https://aiven.io/docs",
        "NewebPay 技術文件與串接規格",
    ]
    return add_bullet_slide(prs, "參考文獻", bullets, page_no, "本簡報內容依專案實作與官方技術文件整理")


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    page = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.1), Inches(11.2), Inches(1.4))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Django + Next.js 電商購物網站專題"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT
    sub = slide.shapes.add_textbox(Inches(2.15), Inches(2.25), Inches(9.0), Inches(0.7))
    p = sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "以前後端分離、MySQL 正式資料庫、購物流程與後台管理為核心"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(17)
    run.font.color.rgb = MUTED
    info = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.05), Inches(3.2), Inches(5.2), Inches(1.25))
    info.fill.solid()
    info.fill.fore_color.rgb = NAVY
    info.line.color.rgb = RGBColor(255, 255, 255)
    tf = info.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "題目：購物網站架設與整合實作\n組員：請填入姓名 / 學號"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    note = slide.shapes.add_textbox(Inches(3.65), Inches(5.05), Inches(6.0), Inches(0.3))
    p = note.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "簡報初稿依專題說明規則與範例章節製作"
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED
    add_footer(slide, page)
    page += 1

    add_bullet_slide(
        prs,
        "目錄",
        [
            "動機與目的",
            "系統架構",
            "開發工具與所需材料",
            "資料庫規劃",
            "流程圖",
            "網站架構與實驗結果",
            "未來方向",
            "參考文獻",
        ],
        page,
    )
    page += 1

    add_section_slide(prs, "動機與目的", page)
    page += 1

    add_two_col_slide(
        prs,
        "動機與目的",
        "專題動機",
        [
            "將課堂的網頁設計、資料庫與後端開發能力整合成完整作品。",
            "不只展示單一頁面，而是完成含會員、商品、訂單與後台的網站。",
            "把專案做成可以部署、可以同步資料、可以延伸維護的形式。",
        ],
        "專題目的",
        [
            "建立前後端分離的電商平台架構。",
            "以 MySQL 作為正式資料來源，完成資料同步與結構管理。",
            "串接金流、物流、媒體儲存與管理後台，接近實務專案流程。",
        ],
        page,
    )
    page += 1

    add_section_slide(prs, "系統架構", page)
    page += 1
    add_system_architecture(prs, page)
    page += 1

    add_section_slide(prs, "開發工具與所需材料", page)
    page += 1
    add_two_col_slide(
        prs,
        "開發工具與所需材料",
        "後端 / 基礎設施",
        [
            "Python 3 + Django 6 + Django REST Framework",
            "MySQL / Aiven 雲端資料庫",
            "Render 部署",
            "Google Cloud Storage 媒體儲存",
            "NewebPay 金流 / 物流",
        ],
        "前端 / 開發工具",
        [
            "Next.js 15 + React 19 + TypeScript 5",
            "VS Code",
            "Git / GitHub",
            "python-pptx 製作簡報初稿",
            "Mermaid / 結構化文件整理網站架構",
        ],
        page,
    )
    page += 1

    add_section_slide(prs, "資料庫規劃", page)
    page += 1
    add_db_slide(prs, page)
    page += 1

    add_section_slide(prs, "流程圖", page)
    page += 1
    add_flow_slide(prs, page)
    page += 1

    add_section_slide(prs, "網站架構與實驗結果", page)
    page += 1
    add_sitemap_slide(prs, page)
    page += 1
    add_results_slide(prs, page)
    page += 1

    add_section_slide(prs, "未來方向", page)
    page += 1
    add_future_slide(prs, page)
    page += 1

    add_section_slide(prs, "參考文獻", page)
    page += 1
    add_reference_slide(prs, page)

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
